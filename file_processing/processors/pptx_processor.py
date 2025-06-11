from io import BytesIO
from pptx import Presentation
import msoffcrypto
import logging
from file_processing.errors import FileProcessingFailedError, FileCorruptionError
from file_processing.file_processor_strategy import FileProcessorStrategy

logger = logging.getLogger(__name__)

class PptxFileProcessor(FileProcessorStrategy):
    """
    Processor for handling PowerPoint (PPTX) files, extracting metadata such as text content,
    author, last modified by, number of slides, and handling encrypted files.

    Attributes:
        metadata (dict): Contains metadata fields such as 'text', 'author', 'last_modified_by',
                         'num_slides', and 'has_password' if the file is opened.
    """

    def __init__(self, file_path: str, open_file: bool = True) -> None:
        """
        Initializes the PptxFileProcessor with the specified file path.

        Args:
            file_path (str): Path to the PPTX file to process.
            open_file (bool): Indicates whether to open and process the file immediately.

        Sets:
            metadata (dict): Populated with a message if `open_file` is False, otherwise initialized with default values.
        """
        super().__init__(file_path, open_file)
        if not open_file:
            logger.debug(f"PPTX file '{self.file_path}' was not opened (open_file=False).")
            self.metadata = {'message': 'File was not opened'}
        else:
            self.metadata = self._default_metadata()

    def _default_metadata(self) -> dict:
        """
        Returns default metadata for an unopened PPTX file.

        Returns:
            dict: Default metadata with 'text', 'author', 'last_modified_by', 'num_slides',
                  and 'has_password' fields.
        """
        return {
            'text': None,
            'author': None,
            'last_modified_by': None,
            'num_slides': None,
            'has_password': False
        }

    def process(self) -> None:
        """
        Extracts metadata from the PPTX file, including text content, author, last modified by,
        number of slides, and checks for encryption.

        Raises:
            FileCorruptionError: If the file is corrupted or encrypted and cannot be opened.
            FileProcessingFailedError: If an error occurs during PPTX file processing.
        """
        logger.info(f"Starting processing of PPTX file '{self.file_path}'.")

        if not self.open_file:
            return

        with open(self.file_path, 'rb') as f:
            file_content = BytesIO(f.read())

        try:
            office_file = msoffcrypto.OfficeFile(file_content)
            if office_file.is_encrypted():
                logger.debug(f"PPTX file '{self.file_path}' is encrypted.")
                self.metadata["has_password"] = True
                return
        except Exception as e:
            logger.error(f"Failed to process PPTX file '{self.file_path}': {e}")
            raise FileCorruptionError(f"File is corrupted: {self.file_path}") from e

        try:
            ppt = Presentation(file_content)
            self.metadata.update({
                'text': self.extract_text_from_pptx(ppt),
                'author': ppt.core_properties.author,
                'last_modified_by': ppt.core_properties.last_modified_by,
                "num_slides": len(ppt.slides),
            })
            logger.info(f"Successfully processed PPTX file '{self.file_path}'.")
        except Exception as e:
            logger.error(f"Failed to process PPTX file '{self.file_path}': {e}")
            raise FileProcessingFailedError(
                f"Error encountered while processing {self.file_path}: {e}"
            )

    def save(self, output_path: str = None) -> None:
        """
        Saves the PPTX file to the specified output path with updated metadata.

        Args:
            output_path (str): Path to save the PPTX file. If None, overwrites the original file.

        Raises:
            FileProcessingFailedError: If an error occurs while saving the PPTX file.
        """
        save_path = output_path or self.file_path
        logger.info(f"Saving PPTX file '{self.file_path}' to '{save_path}'.")

        try:
            ppt = Presentation(self.file_path)

            # Update the core properties (metadata)
            cp = ppt.core_properties
            cp.author = self.metadata.get('author', cp.author)
            cp.last_modified_by = self.metadata.get(
                'last_modified_by', cp.last_modified_by
            )

            ppt.save(save_path)
            logger.info(f"PPTX file '{self.file_path}' saved successfully to '{save_path}'.")
        except Exception as e:
            logger.error(f"Failed to save PPTX file '{self.file_path}' to '{save_path}': {e}")
            raise FileProcessingFailedError(
                f"Error encountered while saving to {save_path}: {e}"
            )

    @staticmethod
    def extract_text_from_pptx(ppt: Presentation) -> str:
        """
        Extracts text from each slide in the PPTX file, including text from tables.

        Args:
            ppt (Presentation): An instance of `Presentation` representing the PPTX file.

        Returns:
            str: Extracted text from the PPTX slides.

        Raises:
            FileProcessingFailedError: If an error occurs while extracting text from the PPTX file.
        """
        try:
            full_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
                    if shape.has_table:
                        for r in shape.table.rows:
                            s = ""
                            for c in r.cells:
                                s += c.text_frame.text + " | "
                            full_text.append(s)
            return '\n'.join(full_text)
        except Exception as e:
            raise FileProcessingFailedError(
                f"Error encountered while extracting text from pptx: {e}"
            )