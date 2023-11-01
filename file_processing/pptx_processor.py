from file_processor_strategy import FileProcessorStrategy
from pptx import Presentation
from zipfile import BadZipFile
from errors import FileProcessingFailedError, FileCorruptionError
import msoffcrypto
from io import BytesIO

class PptxFileProcessor(FileProcessorStrategy):
    def __init__(self, file_path: str, open_file: bool = True) -> None:
        super().__init__(file_path, open_file)
        self.metadata = {'message': 'File was not opened'} if not open_file else self._default_metadata()



    def _default_metadata(self) -> dict:
        return {
            'text': None,
            'author': None,
            'last_modified_by': None,
            'num_slides': None,
            'has_password': False
        }


    def process(self) -> None:
        if not self.open_file:
            return

        with open(self.file_path, 'rb') as f:
            file_content = BytesIO(f.read())

        try:
            office_file = msoffcrypto.OfficeFile(file_content)
            if office_file.is_encrypted():
                self.metadata["has_password"] = True
                return
        except Exception as e:
            raise FileCorruptionError(f"File is corrupted: {self.file_path}")

        try:
            ppt = Presentation(file_content)
            self.metadata.update({'text': self.extract_text_from_pptx(ppt)})
            self.metadata.update({'author': ppt.core_properties.author})
            self.metadata.update({'last_modified_by': ppt.core_properties.last_modified_by})
            self.metadata.update({"num_slides": len(ppt.slides)})
        except Exception as e:
            raise FileProcessingFailedError(f"Error encountered while processing {self.file_path}: {e}")


    def save(self, output_path: str = None) -> None:
        save_path = output_path or self.file_path
        try:
            ppt = Presentation(self.file_path)

            # Update the core properties (metadata)
            cp = ppt.core_properties
            cp.author = self.metadata.get('author', cp.author)
            cp.last_modified_by = self.metadata.get('last_modified_by', cp.last_modified_by)

            ppt.save(save_path)
        except Exception as e:
            raise FileProcessingFailedError(f"Error encountered while saving to {save_path}: {e}")



    @staticmethod
    def extract_text_from_pptx(ppt: Presentation) -> str:
        try:
            full_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
                    if shape.has_table:
                        for r in shape.table.rows:
                            s = ""
                            for c in r.cells:
                                s += c.text_frame.text + " | "
                            full_text.append(s)
            return '\n'.join(full_text)
        except Exception as e:
            raise FileProcessingFailedError(f"Error encountered while extracting text from pptx: {e}")